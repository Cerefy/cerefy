"""Universal text, document, presentation, email, and image parsers.

These parsers extend the πX Cognitive Data Layer so it can ingest a broad range
of enterprise sources without requiring predefined schemas.  Optional
dependencies are imported lazily; if a dependency is missing the parser emits a
warning and returns the best metadata it can extract.
"""

from __future__ import annotations

import email
import re
from io import BytesIO
from pathlib import Path
from typing import Any

import pandas as pd

from app.cognitive_data_layer.parser.base import BaseParser, ParseResult


def _ext(source: str | Path | bytes) -> str:
    if isinstance(source, (str, Path)):
        return Path(source).suffix.lower().lstrip(".")
    return ""


def _read_text(source: str | Path | bytes, options: dict[str, Any] | None) -> tuple[str, list[str]]:
    options = options or {}
    warnings: list[str] = []
    if isinstance(source, (str, Path)):
        with open(source, "rb") as f:
            raw_bytes = f.read()
    else:
        raw_bytes = source
    encoding = options.get("encoding", "utf-8")
    try:
        text = raw_bytes.decode(encoding)
    except UnicodeDecodeError:
        text = raw_bytes.decode("utf-8", errors="replace")
        warnings.append("Fallback decode with replacement characters")
    return text, warnings


def _to_bytesio(source: str | Path | bytes) -> tuple[BytesIO, list[str]]:
    warnings: list[str] = []
    if isinstance(source, (str, Path)):
        with open(source, "rb") as f:
            return BytesIO(f.read()), warnings
    return BytesIO(source), warnings


class TextParser(BaseParser):
    """Parser for plain text files (.txt)."""

    name = "text"
    supported_extensions = ["txt"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        text, warnings = _read_text(source, options)
        df = pd.DataFrame({"line": text.splitlines()})
        return ParseResult(
            raw_data={"sheets": [{"name": "text", "data": df}]},
            format="text",
            metadata={"lines": len(df), "chars": len(text)},
            warnings=warnings,
        )


class MarkdownParser(BaseParser):
    """Parser for Markdown files (.md). Extracts text and headers."""

    name = "markdown"
    supported_extensions = ["md", "markdown"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        text, warnings = _read_text(source, options)
        text = text.replace("\r\n", "\n")
        headers = re.findall(r"^#{1,6}\s+(.+)$", text, flags=re.MULTILINE)
        sections = self._split_markdown_sections(text)
        df = pd.DataFrame({
            "section": [s[0] for s in sections],
            "content": [s[1] for s in sections],
        })
        return ParseResult(
            raw_data={"sheets": [{"name": "markdown", "data": df}]},
            format="markdown",
            metadata={
                "lines": len(text.splitlines()),
                "chars": len(text),
                "headers": headers,
                "sections": len(sections),
            },
            warnings=warnings,
        )

    @staticmethod
    def _split_markdown_sections(text: str) -> list[tuple[str, str]]:
        """Split markdown into (header, content) sections."""
        sections: list[tuple[str, str]] = []
        current_header = ""
        current_lines: list[str] = []
        for line in text.splitlines():
            match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if match:
                if current_lines or current_header:
                    sections.append((current_header, "\n".join(current_lines).strip()))
                current_header = match.group(2).strip()
                current_lines = []
            else:
                current_lines.append(line)
        if current_lines or current_header:
            sections.append((current_header, "\n".join(current_lines).strip()))
        if not sections:
            sections.append(("", text))
        return sections


class HTMLParser(BaseParser):
    """Parser for HTML files (.html, .htm). Extracts visible text and links."""

    name = "html"
    supported_extensions = ["html", "htm"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        text, warnings = _read_text(source, options)
        try:
            from bs4 import BeautifulSoup

            soup = BeautifulSoup(text, "html.parser")
            title = soup.title.get_text(strip=True) if soup.title else ""
            visible_text = soup.get_text(separator="\n", strip=True)
            links = [a.get("href") for a in soup.find_all("a") if a.get("href")]
            df = pd.DataFrame({"text": visible_text.splitlines()})
            return ParseResult(
                raw_data={"sheets": [{"name": "html_text", "data": df}]},
                format="html",
                metadata={
                    "title": title,
                    "links": links,
                    "chars": len(visible_text),
                },
                warnings=warnings,
            )
        except ImportError as exc:  # pragma: no cover
            df = pd.DataFrame({"text": text.splitlines()})
            return ParseResult(
                raw_data={"sheets": [{"name": "html_text", "data": df}]},
                format="html",
                metadata={"chars": len(text)},
                warnings=[f"BeautifulSoup not available: {exc}"],
            )


class RTFParser(BaseParser):
    """Parser for Rich Text Format files (.rtf)."""

    name = "rtf"
    supported_extensions = ["rtf"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        text, warnings = _read_text(source, options)
        try:
            from striprtf.striprtf import rtf_to_text

            plain = rtf_to_text(text)
        except ImportError as exc:  # pragma: no cover
            plain = self._naive_rtf_to_text(text)
            warnings.append(f"striprtf not available; using naive extraction: {exc}")
        df = pd.DataFrame({"line": plain.splitlines()})
        return ParseResult(
            raw_data={"sheets": [{"name": "rtf", "data": df}]},
            format="rtf",
            metadata={"lines": len(df), "chars": len(plain)},
            warnings=warnings,
        )

    @staticmethod
    def _naive_rtf_to_text(text: str) -> str:
        """Fallback: strip RTF control words and braces."""
        text = re.sub(r"\\[a-z]+\d*\s?", "", text)
        text = re.sub(r"[{}]", "", text)
        return text.replace("\\par", "\n")


class PPTXParser(BaseParser):
    """Parser for PowerPoint files (.pptx, .ppt)."""

    name = "pptx"
    supported_extensions = ["pptx", "ppt"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        file, warnings = _to_bytesio(source)
        try:
            from pptx import Presentation

            prs = Presentation(file)
        except ImportError as exc:  # pragma: no cover
            return ParseResult(
                raw_data={"sheets": []},
                format="pptx",
                metadata={},
                warnings=[f"python-pptx not available: {exc}"],
            )

        slides: list[dict[str, Any]] = []
        for i, slide in enumerate(prs.slides, start=1):
            text_lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_lines.append(shape.text)
            slides.append({"slide": i, "text": "\n".join(text_lines)})

        df = pd.DataFrame(slides)
        return ParseResult(
            raw_data={"sheets": [{"name": "slides", "data": df}]},
            format="pptx",
            metadata={"slides": len(prs.slides)},
            warnings=warnings,
        )


class EmailParser(BaseParser):
    """Parser for email files (.eml, .msg)."""

    name = "email"
    supported_extensions = ["eml", "msg"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        text, warnings = _read_text(source, options)
        msg = email.message_from_string(text)
        subject = msg.get("Subject", "")
        sender = msg.get("From", "")
        recipients = msg.get("To", "")
        body = self._extract_email_body(msg)
        df = pd.DataFrame({
            "subject": [subject],
            "from": [sender],
            "to": [recipients],
            "body": [body],
        })
        return ParseResult(
            raw_data={"sheets": [{"name": "email", "data": df}]},
            format="email",
            metadata={
                "subject": subject,
                "from": sender,
                "to": recipients,
                "attachments": [part.get_filename() for part in msg.walk() if part.get_filename()],
            },
            warnings=warnings,
        )

    @staticmethod
    def _extract_email_body(msg: email.message.Message) -> str:
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                payload = part.get_payload(decode=True)
                if isinstance(payload, bytes):
                    return payload.decode("utf-8", errors="replace")
                return str(payload)
        return ""


class ImageParser(BaseParser):
    """Parser for image files (.png, .jpg, .jpeg, .tiff).

    Extracts image metadata. If tesseract is available, also performs OCR.
    """

    name = "image"
    supported_extensions = ["png", "jpg", "jpeg", "tiff", "tif"]

    def can_parse(self, source: str | Path | bytes, hint: str | None = None) -> bool:
        return _ext(source) in self.supported_extensions

    async def parse(
        self,
        source: str | Path | bytes,
        options: dict[str, Any] | None = None,
    ) -> ParseResult:
        file, warnings = _to_bytesio(source)
        try:
            from PIL import Image

            with Image.open(file) as img:
                metadata = {
                    "format": img.format,
                    "mode": img.mode,
                    "width": img.width,
                    "height": img.height,
                }
        except ImportError as exc:  # pragma: no cover
            return ParseResult(
                raw_data={"sheets": []},
                format="image",
                metadata={},
                warnings=[f"Pillow not available: {exc}"],
            )

        ocr_text = ""
        try:
            import pytesseract

            file.seek(0)
            with Image.open(file) as img:
                ocr_text = pytesseract.image_to_string(img)
        except ImportError:
            warnings.append("pytesseract not available; OCR skipped")
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"OCR failed: {exc}")

        df = pd.DataFrame({"text": ocr_text.splitlines() if ocr_text else []})
        return ParseResult(
            raw_data={"sheets": [{"name": "ocr_text", "data": df}]},
            format="image",
            metadata=metadata,
            warnings=warnings,
        )
