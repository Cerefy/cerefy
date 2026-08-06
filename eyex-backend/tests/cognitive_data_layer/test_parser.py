from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pandas as pd
import pytest

from app.cognitive_data_layer.parser import (
    get_parser_registry,
    parse_source,
    register_default_parsers,
)
from app.cognitive_data_layer.parser.plugins import CSVParser, ExcelParser


@pytest.fixture(autouse=True)
def _reset_registry():
    registry = get_parser_registry()
    registry._parsers.clear()
    register_default_parsers()


class TestCSVParser:
    @pytest.mark.asyncio
    async def test_parse_csv(self, tmp_path: Path):
        path = tmp_path / "test.csv"
        path.write_text("Customer,Revenue,Date\nA,100,2026-01-01\nB,200,2026-01-02\n")
        result = await parse_source(path)
        df = result.raw_data["sheets"][0]["data"]
        assert list(df.columns) == ["Customer", "Revenue", "Date"]
        assert len(df) == 2


class TestExcelParser:
    @pytest.mark.asyncio
    async def test_parse_excel(self, tmp_path: Path):
        path = tmp_path / "test.xlsx"
        df = pd.DataFrame({"Client": ["X", "Y"], "Amount": [10, 20]})
        df.to_excel(path, index=False, sheet_name="Sales")
        result = await parse_source(path)
        assert result.format == "excel"
        assert result.raw_data["sheets"][0]["name"] == "Sales"


class TestJSONParser:
    @pytest.mark.asyncio
    async def test_parse_json(self, tmp_path: Path):
        path = tmp_path / "test.json"
        path.write_text('[{"Customer": "A", "Sales": 100}]')
        result = await parse_source(path)
        df = result.raw_data["sheets"][0]["data"]
        assert "Customer" in df.columns


class TestParserRegistry:
    def test_list_parsers(self):
        registry = get_parser_registry()
        names = registry.list_parsers()
        assert "csv" in names
        assert "excel" in names

    def test_csv_can_parse(self):
        parser = CSVParser()
        assert parser.can_parse("data.csv") is True
        assert parser.can_parse("data.xlsx") is False

    def test_excel_can_parse_bytes(self):
        parser = ExcelParser()
        bio = BytesIO()
        pd.DataFrame({"x": [1]}).to_excel(bio, index=False)
        assert parser.can_parse(bio) is True


class TestTextDocumentParsers:
    @pytest.mark.asyncio
    async def test_parse_text(self, tmp_path: Path):
        path = tmp_path / "notes.txt"
        path.write_text("Line one\nLine two\n")
        result = await parse_source(path)
        assert result.format == "text"
        df = result.raw_data["sheets"][0]["data"]
        assert len(df) == 2

    @pytest.mark.asyncio
    async def test_parse_markdown(self, tmp_path: Path):
        from app.cognitive_data_layer.parser.plugins import MarkdownParser

        path = tmp_path / "report.md"
        path.write_text("# Title\n\nSome content.\n## Section\nMore content.")
        result = await MarkdownParser().parse(path)
        assert result.format == "markdown"
        assert "Title" in result.metadata["headers"]

    @pytest.mark.asyncio
    async def test_parse_html(self, tmp_path: Path):
        from app.cognitive_data_layer.parser.plugins import HTMLParser

        path = tmp_path / "page.html"
        path.write_text(
            "<html><head><title>Hello</title></head><body><a href='/x'>Link</a></body></html>"
        )
        result = await HTMLParser().parse(path)
        assert result.format == "html"
        assert result.metadata["title"] == "Hello"
        assert "/x" in result.metadata["links"]

    @pytest.mark.asyncio
    async def test_parse_pptx(self, tmp_path: Path):
        from pptx import Presentation

        from app.cognitive_data_layer.parser.plugins import PPTXParser

        path = tmp_path / "deck.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(0, 0, 100, 50).text_frame.text = "Hello"
        prs.save(path)

        result = await PPTXParser().parse(path)
        assert result.format == "pptx"
        assert result.metadata["slides"] == 1

    @pytest.mark.asyncio
    async def test_parse_email(self, tmp_path: Path):
        from app.cognitive_data_layer.parser.plugins import EmailParser

        path = tmp_path / "message.eml"
        path.write_text(
            "Subject: Test\r\nFrom: a@example.com\r\nTo: b@example.com\r\n\r\nBody text",
            encoding="utf-8",
        )
        result = await EmailParser().parse(path)
        assert result.format == "email"
        assert "Test" in result.metadata["subject"]
        assert "Body text" in result.raw_data["sheets"][0]["data"]["body"].iloc[0]

    @pytest.mark.asyncio
    async def test_parse_image(self, tmp_path: Path):
        from PIL import Image

        from app.cognitive_data_layer.parser.plugins import ImageParser

        path = tmp_path / "image.png"
        Image.new("RGB", (10, 10), color="red").save(path)
        result = await ImageParser().parse(path)
        assert result.format == "image"
        assert result.metadata["width"] == 10
        assert result.metadata["height"] == 10


class TestCodeParser:
    @pytest.mark.asyncio
    async def test_parse_python(self, tmp_path: Path):
        from app.cognitive_data_layer.parser.plugins import CodeParser

        path = tmp_path / "script.py"
        path.write_text(
            "import os\n\nclass Greeter:\n    def greet(self):\n        return 'hi'\n",
            encoding="utf-8",
        )
        result = await CodeParser().parse(path)
        assert result.format == "code"
        assert "Greeter" in result.metadata["classes"]
        assert "greet" in result.metadata["functions"]
        assert "os" in result.metadata["imports"]

    @pytest.mark.asyncio
    async def test_parse_yaml(self, tmp_path: Path):
        from app.cognitive_data_layer.parser.plugins import CodeParser

        path = tmp_path / "config.yaml"
        path.write_text("app:\n  name: test\n", encoding="utf-8")
        result = await CodeParser().parse(path)
        assert result.format == "code"
        assert "app" in result.metadata["top_level_keys"]
